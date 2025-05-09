# utils/text_extractor.py
import requests
from bs4 import BeautifulSoup
import PyPDF2
from pptx import Presentation
import io

def extract_text_from_url(url):
    """Extracts all text content from a given URL."""
    try:
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url
        response = requests.get(url, timeout=15)
        response.raise_for_status()  # Raise an exception for HTTP errors
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove script and style elements
        for script_or_style in soup(["script", "style"]):
            script_or_style.decompose()
        
        # Get text
        text = soup.get_text(separator='\n', strip=True)
        return text
    except requests.exceptions.RequestException as e:
        return f"Error fetching URL: {e}"
    except Exception as e:
        return f"Error parsing URL content: {e}"

def extract_text_from_pdf(file_obj):
    """Extracts text from a PDF file object."""
    try:
        pdf_reader = PyPDF2.PdfReader(file_obj)
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text()
        return text
    except Exception as e:
        return f"Error reading PDF: {e}"

def extract_text_from_ppt(file_obj):
    """Extracts text from a PPTX file object."""
    try:
        prs = Presentation(file_obj)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"Error reading PPTX: {e}"

def extract_text_from_file(uploaded_file):
    """
    Detects file type and extracts text accordingly.
    `uploaded_file` is a Streamlit UploadedFile object.
    """
    if uploaded_file is None:
        return None
        
    file_bytes = io.BytesIO(uploaded_file.getvalue())
    file_type = uploaded_file.type

    if "pdf" in file_type:
        return extract_text_from_pdf(file_bytes)
    elif "vnd.openxmlformats-officedocument.presentationml.presentation" in file_type or uploaded_file.name.endswith('.pptx'): # PPTX
        return extract_text_from_ppt(file_bytes)
    else:
        return "Unsupported file type for text extraction."